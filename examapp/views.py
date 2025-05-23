# views.py
from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
from django.views.decorators.http import require_POST
from django.contrib import messages
from .utils import *
from django.http import FileResponse, HttpResponse, JsonResponse
from django.db.models import Q
from pptx import Presentation
from pptx.util import Inches, Pt
from .models import *
import tempfile
import os 
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from django.utils import timezone
from .decorators import *
from .forms import *
from .authentication import EmailOrPhoneBackend
from django.contrib.auth import authenticate, login, logout
from django.utils.http import unquote



def home(request):
    return render(request, 'home.html')

@staff_member_required
@require_POST
def undo_last_exam_action(request):
    exam_ids = request.session.get('undo_exam_ids', [])

    if exam_ids:
        Exam.objects.filter(id__in=exam_ids).delete()
        messages.success(request, "✅ Undo successful! Exams deleted.")
        request.session.pop('undo_exam_ids')
    else:
        messages.warning(request, "⚠️ No recent exams to undo.")

    return redirect('create_exam')


@login_required(login_url='login')
@staff_member_required
def create_exam_page(request):
    if request.method == 'POST':
        try:
            number = int(request.POST.get("number", 0))
            if number <= 0:
                raise ValueError("Number must be greater than 0")
            
            exams_created, created_exam_ids = auto_create_exams(number)
            request.session['undo_exam_ids'] = created_exam_ids
            request.session['show_undo'] = True  # Add flag

            messages.success(request, f"{exams_created} exam(s) created successfully!")
            return redirect('create_exam')
        except (ValueError, TypeError):
            messages.error(request, "Invalid number of exams.")

    # Show last 10 Ibivanze exams
    ibivanze_type = ExamType.objects.filter(name='Ibivanze').first()
    recent_exams = Exam.objects.filter(exam_type=ibivanze_type).order_by('-created_at')[:10]
    
    context = {
        'recent_exams': recent_exams,
        'show_undo': request.session.pop('show_undo', False),
        'has_undo_ids': bool(request.session.get('undo_exam_ids')),
    }
    return render(request, 'exams/create_exam.html', context)


def generate_exam_ppt(exam, bg_color, text_color, highlight_color, correct_choice_text):
    prs = Presentation()
    prs.slide_width = Inches(14.5)
    prs.slide_height = Inches(8.5)

    blank_layout = next(layout for layout in prs.slide_layouts if layout.name == "Blank")
    questions = list(exam.questions.all())[:20]

    for index, question in enumerate(questions, start=1):
        # --- Question Slide ---
        slide_q = prs.slides.add_slide(blank_layout)
        background = slide_q.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        tx_box = slide_q.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{index}. {question.question_text}"
        p.font.size = Pt(40)
        p.font.color.rgb = text_color

        # Question Image
        if question.question_sign:
            slide_q.shapes.add_picture(question.question_sign.sign_image.path, Inches(2), Inches(1.5), height=Inches(2.5))

        # Choices
        choice_box = slide_q.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4))
        tf = choice_box.text_frame
        tf.word_wrap = True
        tf.auto_size = False
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)

        image_top = Inches(5)
        image_left = Inches(0.8)
        image_width = Inches(2.5)
        image_height = Inches(2.0)
        img_per_row = 3
        img_count = 0

        for choice in question.get_choices():
            content = choice["content"]
            is_img = choice["type"] == "image"
            letter = chr(96 + choice["id"])

            if is_img:
                left = image_left + (img_count % img_per_row) * (image_width + Inches(0.3))
                top = image_top + (img_count // img_per_row) * (image_height + Inches(0.3))
                img_path = content.replace('/media/', 'media/')
                slide_q.shapes.add_picture(img_path, left, top, width=image_width, height=image_height)
                img_count += 1
            else:
                p = tf.add_paragraph()
                p.text = f"{letter}) {content}"
                p.font.size = Pt(38)
                p.font.color.rgb = text_color

        # --- Answer Slide ---
        slide_ans = prs.slides.add_slide(blank_layout)
        fill = slide_ans.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        correct = next((c for c in question.get_choices() if c["is_correct"]), None)
        if correct:
            content = correct["content"]
            letter = chr(96 + correct["id"])
            ans_text = f"({letter}) {content if correct['type'] == 'text' else 'Ifoto'}"

            box = slide_ans.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10), Inches(1.2))
            tf = box.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = ans_text
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = correct_choice_text
            p.alignment = PP_ALIGN.CENTER
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = highlight_color

            if correct["type"] == "image":
                slide_ans.shapes.add_picture(
                    correct["content"].replace('/media/', 'media/'),
                    Inches(2.5), Inches(4.5), height=Inches(2.5)
                )

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


@login_required(login_url='login')
@staff_member_required
def pptx(request):
    return render(request, 'exams/pptx.html')



def hex_to_rgbcolor(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

@login_required(login_url='login')
@staff_member_required
def exam_ppt_preview(request):
    user = request.user
    exam = Exam.objects.latest("created_at")
    first_question = exam.questions.first()

    if request.method == "POST":
        bg_color = request.POST.get("bg_color", "#ffffff")
        text_color = request.POST.get("text_color", "#000000")
        highlight_color = request.POST.get("highlight_color", "#00ff00")
        correct_choice_text = request.POST.get("correct_choice_text", "#000000")
        save_colors = request.POST.get("save_defaults") == "on"

        if save_colors:
            user.ppt_bg_color = bg_color
            user.ppt_text_color = text_color
            user.ppt_highlight_color = highlight_color
            user.ppt_correct_choice_text = correct_choice_text
            user.save(update_fields=["ppt_bg_color", "ppt_text_color", "ppt_highlight_color", "ppt_correct_choice_text"])

        ppt_file_path = generate_exam_ppt(
            exam,
            hex_to_rgbcolor(bg_color),
            hex_to_rgbcolor(text_color),
            hex_to_rgbcolor(highlight_color),
            hex_to_rgbcolor(correct_choice_text)
        )

        with open(ppt_file_path, 'rb') as f:
            response = HttpResponse(
                f.read(),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = f'attachment; filename="Exam_{timezone.now().date()}.pptx"'
        os.remove(ppt_file_path)
        return response

    # Pre-fill with saved user values
    context = {
        'exam': exam,
        'question': first_question,
        'bg_color': user.ppt_bg_color,
        'text_color': user.ppt_text_color,
        'highlight_color': user.ppt_highlight_color,
        'correct_choice_text': user.ppt_correct_choice_text,
    }
    return render(request, 'exams/pptx_preview.html', context)


def check_unique_field(request):
    field = request.GET.get("field")
    value = request.GET.get("value")
    response = {"exists": False}

    if field == "email" and value:
        response["exists"] = User.objects.filter(email__iexact=value).exists()
    elif field == "phone_number" and value:
        
        if len(value) >= 10:
            response["exists"] = User.objects.filter(phone_number__icontains=value).exists()
        # response["exists"] = User.objects.filter(phone_number__contains=value).exists()
    elif field == "name" and value:
        response["exists"] = User.objects.filter(name__iexact=value).exists()

    return JsonResponse(response)


@redirect_authenticated_users
def register_view(request):
    if request.method == 'POST':
        form = RegisterForm(request.POST)
        if form.is_valid():
            user = form.save(commit=False)
            user.set_password(form.cleaned_data["password1"])
            user.save()

            # Send OTP to the user's phone number
            if form.cleaned_data.get("phone_number"):
                user.otp_verified = True  
                user.save()
                messages.success(request, 'Kwiyandikisha muri Kigali Driving School byagenze neza')
                login(request, user, backend='django.contrib.auth.backends.ModelBackend')
                return redirect('home')            

            
    else:
        form = RegisterForm()
    return render(request, 'registration/register.html', {'form': form})


@redirect_authenticated_users
def login_view(request):
    page='login'
    
    # show_modal = request.GET.get('login') is not None
    if request.method == "POST":
        form = LoginForm(request.POST)
        
        if form.is_valid():
            username = form.cleaned_data["username"]
            password = form.cleaned_data["password"]

            # Fetch user by email or phone number
             # Normalize phone number if needed
            if "@" not in username:  
                username = EmailOrPhoneBackend().normalize_phone_number(username)

            # Fetch user by email or phone number
            user = UserProfile.objects.filter(Q(phone_number=username) | Q(email=username)).first()
            if user:
                # Ensure phone_number is set to None if empty
                if user.phone_number == "":
                    user.phone_number = None
                    user.save(update_fields=["phone_number"])
                
                if user.email and not user.otp_verified:
                    messages.error(request, "Please banza wuzuze kode yoherejwe yemeza ko email ari yawe.")
                    return redirect("verify_otp", user_id=user.id)
                authenticated_user = authenticate(request, username=username, password=password)
                
                if authenticated_user:
                    login(request, authenticated_user)
                    messages.success(request, "Kwinjira bikozwe neza cyane! Ikaze nanone.")
                    return redirect("home")
                else:
                    messages.error(request, "Ijambobanga ritariryo, ongera ugerageze.")
            else:
                messages.error(request, "Iyi konti ntago ibaho!!!")

        # Handle form validation errors
        for field, errors in form.errors.items():
            for error in errors:
                messages.error(request, f"{field.capitalize()}: {error}")

    else:
        form = LoginForm()
        
    context = {
        "form": form,
        "page":page
    }
    return render(request, "base.html",context)


@require_POST
@login_required(login_url='login')
def user_logout(request):
    logout(request)
    messages.info(request, "Gusohoka byakunze.")
    return redirect('login')

